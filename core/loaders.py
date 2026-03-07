"""
core/loaders.py
Handles dynamic file loading and text chunking for the RAG pipeline.
Supports: PDF, DOCX, PPTX, XLSX, CSV
"""

import os
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter


def load_file_dynamically(file_path: str) -> list[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    loaders = {
        ".pdf":  _load_pdf,
        ".docx": _load_docx,
        ".pptx": _load_pptx,
        ".xlsx": _load_xlsx,
        ".csv":  _load_csv,
    }
    if ext not in loaders:
        raise ValueError(f"Unsupported file type: {ext}")
    docs = loaders[ext](file_path)
    if not docs:
        raise ValueError(f"No text could be extracted from {os.path.basename(file_path)}. The file may be empty or image-based.")
    return docs


def get_chunks(docs: list[Document], chunk_size: int = 800, chunk_overlap: int = 100) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_documents(docs)


def _load_pdf(path: str) -> list[Document]:
    from langchain_community.document_loaders import PyMuPDFLoader
    return PyMuPDFLoader(path).load()


def _load_docx(path: str) -> list[Document]:
    from langchain_community.document_loaders import Docx2txtLoader
    return Docx2txtLoader(path).load()


def _load_pptx(path: str) -> list[Document]:
    from pptx import Presentation
    prs = Presentation(path)
    docs = []
    for i, slide in enumerate(prs.slides):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"source": path, "slide": i + 1}
            ))
    return docs


def _load_xlsx(path: str) -> list[Document]:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    docs = []
    for sheet in wb.sheetnames:
        rows = [
            " | ".join(str(c) for c in row if c is not None)
            for row in wb[sheet].iter_rows(values_only=True)
        ]
        rows = [r for r in rows if r.strip()]
        if rows:
            docs.append(Document(
                page_content="\n".join(rows),
                metadata={"source": path, "sheet": sheet}
            ))
    return docs


def _load_csv(path: str) -> list[Document]:
    import csv
    with open(path, newline="", encoding="utf-8-sig") as f:
        rows = [
            " | ".join(c.strip() for c in row if c.strip())
            for row in csv.reader(f)
        ]
    rows = [r for r in rows if r]
    return [Document(page_content="\n".join(rows), metadata={"source": path})]